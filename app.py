import os
import re
import time
from pathlib import Path

import pdfplumber
import streamlit as st
from docx import Document
from dotenv import load_dotenv
from google import genai
from openai import OpenAI
from pptx import Presentation

load_dotenv()

st.set_page_config(
    page_title="EduMentor AI",
    page_icon="🎓",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
.stApp { background: linear-gradient(135deg, #F8FAFC 0%, #EEF2FF 45%, #ECFDF5 100%); }
.main .block-container { padding-top: 2rem; padding-bottom: 2rem; max-width: 1300px; }
[data-testid="stSidebar"] { background: linear-gradient(180deg, #0F172A 0%, #1E293B 100%); }
[data-testid="stSidebar"] * { color: #F8FAFC !important; }

.hero-title { color:#0F172A; font-size:58px; font-weight:850; line-height:1.08; margin-bottom:16px; }
.hero-subtitle { color:#334155; font-size:21px; line-height:1.6; max-width:900px; margin-bottom:32px; }
.section-title { color:#0F172A; font-size:34px; font-weight:800; margin-top:36px; margin-bottom:20px; }

.premium-card {
    background:rgba(255,255,255,0.92);
    border:1px solid rgba(15,23,42,0.08);
    border-radius:24px;
    padding:28px;
    min-height:170px;
    box-shadow:0 18px 45px rgba(15,23,42,0.08);
}
.card-title { color:#0F172A; font-size:23px; font-weight:800; margin-bottom:12px; }
.card-text { color:#475569; font-size:16px; line-height:1.7; }

.upload-card {
    background:#FFFFFF;
    border:2px dashed rgba(16,185,129,0.45);
    border-radius:26px;
    padding:32px;
    margin-bottom:22px;
    box-shadow:0 16px 40px rgba(15,23,42,0.08);
}
.upload-title { color:#0F172A; font-size:26px; font-weight:800; margin-bottom:10px; }
.upload-text { color:#475569; font-size:17px; line-height:1.6; }

.preview-box {
    background:#FFFFFF;
    border:1px solid rgba(15,23,42,0.10);
    border-radius:20px;
    padding:24px;
    max-height:420px;
    overflow-y:auto;
    box-shadow:0 14px 35px rgba(15,23,42,0.08);
    color:#0F172A;
    white-space:pre-wrap;
    line-height:1.7;
}

[data-testid="stAlert"] { border-radius:16px; border:1px solid rgba(16,185,129,0.35); }
[data-testid="stAlert"] * { color:#0F172A !important; }

.stButton > button {
    background:linear-gradient(135deg,#2563EB,#10B981);
    color:white;
    border:none;
    border-radius:14px;
    padding:0.75rem 1.3rem;
    font-weight:700;
}
[data-testid="stFileUploader"] {
    background:#FFFFFF;
    border-radius:18px;
    padding:10px;
    border:1px solid rgba(15,23,42,0.08);
}
header { background:transparent !important; }
</style>
""", unsafe_allow_html=True)

UPLOAD_FOLDER = Path("data/uploads")

PRIMARY_OPENAI_MODEL = "gpt-4.1-mini"
FALLBACK_GEMINI_MODEL = "gemini-2.5-flash"
MAX_RETRIES = 2

PROMPTS = {
    "explain": """
You are EduMentor AI, an expert AI tutor for Indian university students.

Create a clear, beginner-friendly explanation.

Use proper Markdown formatting.

Format:
# 📘 Topic Overview
# 🧠 Simple Explanation
# 🔑 Important Concepts

For each concept:
## Concept Name
**Meaning:**  
Explain simply.

**Why it is important:**  
Explain importance.

**Example:**  
Give a real-life or technical example.

# 📝 Exam Tips
# ⚡ Quick Revision Points
""",

    "summary": """
Create a clean exam-focused summary using proper Markdown.

Format:
# 📝 Short Summary
# 📚 Detailed Summary
# 🔑 Key Points
# 📌 Important Terms
# ⚡ Quick Revision Notes
""",

    "mcq": """
Generate 15 MCQs from the uploaded notes.

Use this exact Markdown format:

# ❓ MCQs

## Question 1

**Question:**  
Write the question here.

**Options:**  
A) Option 1  
B) Option 2  
C) Option 3  
D) Option 4  

**Correct Answer:**  
A/B/C/D

**Explanation:**  
Explain why the answer is correct.

Repeat for all questions.
""",

    "flashcards": """
Create flashcards from the uploaded notes.

Use this exact Markdown format:

# 🧠 Flashcards

## Flashcard 1

**Question:**  
Write the question here.

**Answer:**  
Write the answer clearly here.

**Key Point:**  
Mention the most important point to remember.

Repeat for all flashcards.
""",

    "viva": """
Generate viva questions and answers from the uploaded notes.

Use this exact Markdown format:

# 🎤 Viva Questions

## Question 1

**Question:**  
Write the viva question here.

**Answer:**  
Write the answer clearly below the question.

**Important Point:**  
Mention what the student should remember.

Repeat for all viva questions.

# 📌 Important Viva Topics
List important topics for oral exam preparation.
"""
}


def get_unique_file_path(file_name):
    UPLOAD_FOLDER.mkdir(parents=True, exist_ok=True)
    original_path = UPLOAD_FOLDER / file_name

    if not original_path.exists():
        return original_path

    stem = original_path.stem
    suffix = original_path.suffix
    counter = 1

    while True:
        new_path = UPLOAD_FOLDER / f"{stem}_{counter}{suffix}"
        if not new_path.exists():
            return new_path
        counter += 1


def save_uploaded_file(uploaded_file):
    file_path = get_unique_file_path(uploaded_file.name)
    with open(file_path, "wb") as file:
        file.write(uploaded_file.getbuffer())
    return file_path


def clean_text(text):
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(line.strip() for line in text.split("\n"))
    return text.strip()


def extract_text_from_pdf(file_path):
    extracted_text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                extracted_text += text + "\n\n"
    return extracted_text.strip()


def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read().strip()


def extract_text_from_docx(file_path):
    document = Document(file_path)
    extracted_text = ""
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            extracted_text += paragraph.text + "\n\n"
    return extracted_text.strip()


def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    extracted_text = ""

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"

        if slide_text.strip():
            extracted_text += f"Slide {slide_number}:\n{slide_text}\n\n"

    return extracted_text.strip()


def process_document(file_path):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif suffix == ".txt":
        text = extract_text_from_txt(file_path)
    elif suffix == ".docx":
        text = extract_text_from_docx(file_path)
    elif suffix == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        text = ""

    return clean_text(text)


def build_prompt(mode, text):
    safe_text = text[:10000]
    return f"""
{PROMPTS[mode]}

Uploaded Notes:
{safe_text}
"""


def call_openai(prompt):
    api_key = os.getenv("OPENAI_API_KEY")

    if not api_key:
        raise ValueError("OpenAI API key not found.")

    client = OpenAI(api_key=api_key)

    response = client.responses.create(
        model=PRIMARY_OPENAI_MODEL,
        input=prompt,
    )

    return response.output_text


def call_gemini(prompt):
    api_key = os.getenv("GEMINI_API_KEY")

    if not api_key:
        raise ValueError("Gemini API key not found.")

    client = genai.Client(api_key=api_key)

    response = client.models.generate_content(
        model=FALLBACK_GEMINI_MODEL,
        contents=prompt
    )

    return response.text


def generate_ai_response(mode, text):
    prompt = build_prompt(mode, text)

    for _ in range(MAX_RETRIES):
        try:
            response_text = call_openai(prompt)

            if response_text and response_text.strip():
                return True, response_text

        except Exception:
            time.sleep(1)

    try:
        response_text = call_gemini(prompt)

        if response_text and response_text.strip():
            return True, response_text

    except Exception:
        pass

    return False, (
        "⚠️ **AI service is currently busy.**\n\n"
        "Please wait **15–30 seconds** and try again.\n\n"
        "✅ Your uploaded notes are safe."
    )


def render_ai_page(title, subtitle, mode, button_text):
    st.markdown(f'<div class="hero-title">{title}</div>', unsafe_allow_html=True)

    st.markdown(
        f"""
        <div class="hero-subtitle">
            {subtitle}
        </div>
        """,
        unsafe_allow_html=True
    )

    if not st.session_state["extracted_text"]:
        st.warning("Please upload and process a document first from the Upload Notes page.")
        return

    st.markdown('<div class="section-title">Uploaded Document</div>', unsafe_allow_html=True)

    st.markdown(f"""
    <div class="premium-card">
        <div class="card-title">📄 {st.session_state["uploaded_file_name"]}</div>
        <div class="card-text">Your document is processed and ready.</div>
    </div>
    """, unsafe_allow_html=True)

    if st.button(button_text):
        with st.spinner("EduMentor AI is generating your result..."):
            success, response = generate_ai_response(
                mode,
                st.session_state["extracted_text"]
            )

        if success:
            st.session_state["ai_outputs"][mode] = response
            st.session_state["ai_errors"].pop(mode, None)
        else:
            st.session_state["ai_errors"][mode] = response

    if st.session_state["ai_errors"].get(mode):
        st.warning(st.session_state["ai_errors"][mode])

    if st.session_state["ai_outputs"].get(mode):
        st.markdown('<div class="section-title">Generated Output</div>', unsafe_allow_html=True)
        st.markdown(st.session_state["ai_outputs"][mode])


if "extracted_text" not in st.session_state:
    st.session_state["extracted_text"] = ""

if "uploaded_file_name" not in st.session_state:
    st.session_state["uploaded_file_name"] = ""

if "uploaded_file_signature" not in st.session_state:
    st.session_state["uploaded_file_signature"] = ""

if "ai_outputs" not in st.session_state:
    st.session_state["ai_outputs"] = {}

if "ai_errors" not in st.session_state:
    st.session_state["ai_errors"] = {}


st.sidebar.markdown("## 🎓 EduMentor AI")
st.sidebar.markdown("AI Learning Assistant")
st.sidebar.divider()

page = st.sidebar.radio(
    "Navigation",
    [
        "🏠 Dashboard",
        "📂 Upload Notes",
        "🤖 AI Explanation",
        "📝 Summary",
        "❓ MCQs",
        "🧠 Flashcards",
        "🎤 Viva",
        "📚 Resources",
        "⚙ Settings"
    ],
    key="navigation"
)


if page == "🏠 Dashboard":
    st.markdown(
        """
        <div class="hero-title">
            Learn Smarter.<br>
            Understand Faster.
        </div>
        """,
        unsafe_allow_html=True
    )

    st.markdown(
        """
        <div class="hero-subtitle">
            Upload your university notes and generate explanations, summaries, MCQs, flashcards, and viva questions.
        </div>
        """,
        unsafe_allow_html=True
    )

    st.markdown('<div class="section-title">Core Features</div>', unsafe_allow_html=True)

    col1, col2, col3 = st.columns(3)

    with col1:
        st.markdown("""
        <div class="premium-card">
            <div class="card-title">🤖 AI Explanation</div>
            <div class="card-text">Understand difficult topics in beginner-friendly language.</div>
        </div>
        """, unsafe_allow_html=True)

    with col2:
        st.markdown("""
        <div class="premium-card">
            <div class="card-title">📝 Smart Summary</div>
            <div class="card-text">Convert long notes into summaries and revision points.</div>
        </div>
        """, unsafe_allow_html=True)

    with col3:
        st.markdown("""
        <div class="premium-card">
            <div class="card-title">❓ Exam Practice</div>
            <div class="card-text">Generate MCQs, flashcards, and viva questions.</div>
        </div>
        """, unsafe_allow_html=True)


elif page == "📂 Upload Notes":
    st.markdown('<div class="hero-title">Upload Notes</div>', unsafe_allow_html=True)

    st.markdown(
        """
        <div class="hero-subtitle">
            Upload PDFs, DOCX files, PPTs, TXT files, or images.
        </div>
        """,
        unsafe_allow_html=True
    )

    st.markdown("""
    <div class="upload-card">
        <div class="upload-title">☁️ Drag & Drop Study Material</div>
        <div class="upload-text">Supported formats: PDF • DOCX • PPTX • TXT • PNG • JPG • JPEG</div>
    </div>
    """, unsafe_allow_html=True)

    uploaded_file = st.file_uploader(
        "Choose your study material",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg"]
    )

    if uploaded_file is not None:
        file_signature = f"{uploaded_file.name}_{uploaded_file.size}"
        saved_path = save_uploaded_file(uploaded_file)

        st.success("✅ File uploaded and saved successfully!")

        with st.spinner("Processing and cleaning document..."):
            extracted_text = process_document(saved_path)

        if extracted_text:
            if st.session_state["uploaded_file_signature"] != file_signature:
                st.session_state["ai_outputs"] = {}
                st.session_state["ai_errors"] = {}

            st.session_state["extracted_text"] = extracted_text
            st.session_state["uploaded_file_name"] = uploaded_file.name
            st.session_state["uploaded_file_signature"] = file_signature

            st.markdown('<div class="section-title">Cleaned Text Preview</div>', unsafe_allow_html=True)

            st.markdown(
                f"""
                <div class="preview-box">
                {extracted_text[:4000]}
                </div>
                """,
                unsafe_allow_html=True
            )

            st.success("✅ Document processed and ready for all AI tools!")

        elif saved_path.suffix.lower() in [".pdf", ".txt", ".docx", ".pptx"]:
            st.warning("No readable text was found in this file.")

        else:
            st.info("Text extraction for images will be added in the OCR module.")

    else:
        st.info("Please upload a file to continue.")


elif page == "🤖 AI Explanation":
    render_ai_page(
        "AI Explanation",
        "Explain uploaded notes in simple, beginner-friendly language.",
        "explain",
        "✨ Generate Explanation"
    )


elif page == "📝 Summary":
    render_ai_page(
        "Smart Summary",
        "Generate clean summaries, key points, and revision notes.",
        "summary",
        "📝 Generate Summary"
    )


elif page == "❓ MCQs":
    render_ai_page(
        "MCQ Generator",
        "Generate exam-style MCQs from uploaded notes.",
        "mcq",
        "❓ Generate MCQs"
    )


elif page == "🧠 Flashcards":
    render_ai_page(
        "Flashcards",
        "Create question-answer flashcards for quick revision.",
        "flashcards",
        "🧠 Generate Flashcards"
    )


elif page == "🎤 Viva":
    render_ai_page(
        "Viva Questions",
        "Generate viva questions with answers for oral exams.",
        "viva",
        "🎤 Generate Viva Questions"
    )


else:
    st.markdown(f'<div class="hero-title">{page}</div>', unsafe_allow_html=True)
    st.markdown(
        """
        <div class="hero-subtitle">
            This module will be developed in upcoming lessons.
        </div>
        """,
        unsafe_allow_html=True
    )
    st.info("🚧 Module under development.")