import re
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation


UPLOAD_FOLDER = Path("data/uploads")


def get_unique_file_path(file_name):
    UPLOAD_FOLDER.mkdir(parents=True, exist_ok=True)

    original_path = UPLOAD_FOLDER / file_name

    if not original_path.exists():
        return original_path

    stem = original_path.stem
    suffix = original_path.suffix
    counter = 1

    while True:
        new_path = UPLOAD_FOLDER / f"{stem}_{counter}{suffix}"

        if not new_path.exists():
            return new_path

        counter += 1


def save_uploaded_file(uploaded_file):
    file_path = get_unique_file_path(uploaded_file.name)

    with open(file_path, "wb") as file:
        file.write(uploaded_file.getbuffer())

    return file_path


def clean_text(text):
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(line.strip() for line in text.split("\n"))

    return text.strip()


def extract_text_from_pdf(file_path):
    extracted_text = ""

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()

            if text:
                extracted_text += text + "\n\n"

    return extracted_text.strip()


def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read().strip()


def extract_text_from_docx(file_path):
    document = Document(file_path)
    extracted_text = ""

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            extracted_text += paragraph.text + "\n\n"

    return extracted_text.strip()


def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    extracted_text = ""

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"

        if slide_text.strip():
            extracted_text += f"Slide {slide_number}:\n{slide_text}\n\n"

    return extracted_text.strip()


def process_document(file_path):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif suffix == ".txt":
        text = extract_text_from_txt(file_path)
    elif suffix == ".docx":
        text = extract_text_from_docx(file_path)
    elif suffix == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        text = ""

    return clean_text(text)